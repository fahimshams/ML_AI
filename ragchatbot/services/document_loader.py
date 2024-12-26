import os
from pptx import Presentation

def load_documents_from_directory(directory_path):
    print("==== Loading documents from directory ====")
    documents = []
    for filename in os.listdir(directory_path):
        if filename.endswith(".pptx"):
            pptx_path = os.path.join(directory_path, filename)
            presentation = Presentation(pptx_path)
            text = ""
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            documents.append({"id": filename, "text": text})
    return documents